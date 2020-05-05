from argparse import ArgumentParser
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

parser = ArgumentParser(description='Extract images from powerpoint presentations')
parser.add_argument('file', metavar='f', nargs='+', type=str, help='The file path to extract images')
args = parser.parse_args()

# Global counter
n = 0
written_shas = {}


def write_image(shape):
    global n
    image = shape.image
    image_bytes = image.blob
    image_filename = '/Users/overlord/extracted-images/image{:03d}.{}'.format(n, image.ext)
    written_shas[image.sha1] = True
    n += 1
    print(image_filename)
    print(image.filename)
    with open(image_filename, 'wb') as f:
        f.write(image_bytes)


def visitor(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            visitor(s)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.image.sha1 not in written_shas:
        write_image(shape)


def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            visitor(shape)


print(args.file)
iter_picture_shapes(Presentation(args.file[0]))
